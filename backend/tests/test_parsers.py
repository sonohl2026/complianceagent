import io

import docx
import fitz
import openpyxl
import pytest
from pptx import Presentation
from pptx.util import Inches

from app.services.parsing.base import ParsingError
from app.services.parsing.dispatch import parse_document
from app.services.parsing.parsers.csv_parser import parse_csv
from app.services.parsing.parsers.html_parser import parse_html
from app.services.parsing.parsers.text_parser import parse_markdown, parse_plain_text


def _build_pdf_bytes(text: str) -> bytes:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), text, fontsize=12)
    data = doc.tobytes()
    doc.close()
    return data


def _build_docx_bytes() -> bytes:
    document = docx.Document()
    document.add_heading("Intended Use", level=1)
    document.add_paragraph("SonoHL is an investigational acoustic-sensing platform.")
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _build_pptx_bytes() -> bytes:
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = "Product Overview"
    body = slide.placeholders[1]
    body.text_frame.text = "Multi-point acoustic acquisition"
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def _build_xlsx_bytes() -> bytes:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "CodeMatrix"
    ws.append(["Code", "System", "Status"])
    ws.append(["0XXXT", "CPT Category III", "EXPERT REVIEW REQUIRED"])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def test_parse_pdf_extracts_text_and_page_number():
    parsed = parse_document(".pdf", _build_pdf_bytes("Investigational device notice"))
    assert any("Investigational device notice" in b.text for b in parsed.blocks)
    assert parsed.blocks[0].page_number == 1


def test_parse_pdf_rejects_scanned_no_text_layer_pdf():
    doc = fitz.open()
    doc.new_page()  # blank page, no text layer at all
    data = doc.tobytes()
    doc.close()
    with pytest.raises(ParsingError, match="no extractable text layer"):
        parse_document(".pdf", data)


def test_parse_docx_captures_heading_and_body():
    parsed = parse_document(".docx", _build_docx_bytes())
    headings = [b for b in parsed.blocks if b.heading_level > 0]
    body = [b for b in parsed.blocks if b.heading_level == 0]
    assert headings and headings[0].text == "Intended Use"
    assert any("investigational acoustic-sensing platform" in b.text for b in body)


def test_parse_pptx_captures_slide_title_and_number():
    parsed = parse_document(".pptx", _build_pptx_bytes())
    assert parsed.title == "Product Overview"
    assert all(b.page_number == 1 for b in parsed.blocks)
    assert any("Multi-point acoustic acquisition" in b.text for b in parsed.blocks)


def test_parse_xlsx_captures_rows_and_sheet_name():
    parsed = parse_document(".xlsx", _build_xlsx_bytes())
    assert any(b.text == "CodeMatrix" for b in parsed.blocks)
    assert any("0XXXT" in b.text for b in parsed.blocks)


def test_parse_csv_first_row_is_header():
    content = b"Code,System,Status\n0XXXT,CPT Category III,EXPERT REVIEW REQUIRED\n"
    parsed = parse_csv(content)
    assert parsed.blocks[0].heading_level == 1
    assert parsed.blocks[0].text == "Code | System | Status"


def test_parse_html_extracts_headings_and_paragraphs():
    html = b"<html><head><title>SonoHL</title></head><body><h1>Overview</h1><p>Investigational device.</p></body></html>"
    parsed = parse_html(html)
    assert parsed.title == "SonoHL"
    assert parsed.blocks[0].heading_level == 1
    assert parsed.blocks[0].text == "Overview"
    assert parsed.blocks[1].text == "Investigational device."


def test_parse_html_rejects_empty_document():
    with pytest.raises(ParsingError):
        parse_html(b"<html><body></body></html>")


def test_parse_markdown_headings_and_title():
    md = b"# SonoHL Overview\n\nSonoHL is investigational.\n\n## Intended Use\n\nClinician support only.\n"
    parsed = parse_markdown(md)
    assert parsed.title == "SonoHL Overview"
    heading_texts = [b.text for b in parsed.blocks if b.heading_level > 0]
    assert heading_texts == ["SonoHL Overview", "Intended Use"]


def test_parse_plain_text_splits_on_blank_lines():
    parsed = parse_plain_text(b"Paragraph one.\n\nParagraph two.")
    assert len(parsed.blocks) == 2
    assert parsed.blocks[0].text == "Paragraph one."

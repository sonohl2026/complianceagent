import io

from pptx import Presentation

from app.services.parsing.base import Block, ParsedDocument, ParsingError


def parse_pptx(content: bytes) -> ParsedDocument:
    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as exc:
        raise ParsingError(f"Failed to open PPTX: {exc}") from exc

    blocks: list[Block] = []
    title = None

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            is_title_placeholder = getattr(shape, "is_placeholder", False) and getattr(
                shape.placeholder_format, "idx", None
            ) == 0
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip() or paragraph.text.strip()
                if not text:
                    continue
                heading_level = 1 if is_title_placeholder else 0
                if is_title_placeholder and title is None:
                    title = text
                blocks.append(Block(text=text, page_number=slide_index, heading_level=heading_level))

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            blocks.append(
                Block(
                    text=f"[Speaker notes] {slide.notes_slide.notes_text_frame.text.strip()}",
                    page_number=slide_index,
                    heading_level=0,
                )
            )

    if not blocks:
        raise ParsingError("PPTX contains no extractable text")

    return ParsedDocument(title=title, blocks=blocks)
